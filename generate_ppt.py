#!/usr/bin/env python3
"""
PPT Generator - Generate proper PPTX slides with text + AI illustrations.

Uses python-pptx for slide structure/text and Qwen wanx for illustrations.
"""

import argparse
import json
import os
import sys
from datetime import datetime
from pathlib import Path
from typing import Optional

from dotenv import load_dotenv
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# =============================================================================
# Constants
# =============================================================================

OUTPUT_BASE_DIR = "outputs"

# Dark theme colors
COLOR_BG        = RGBColor(0x0D, 0x1B, 0x2A)  # deep navy
COLOR_ACCENT    = RGBColor(0x00, 0xB4, 0xD8)  # cyan
COLOR_HEADING   = RGBColor(0xFF, 0xFF, 0xFF)  # white
COLOR_BODY      = RGBColor(0xCA, 0xD3, 0xE0)  # light grey
COLOR_COVER_SUB = RGBColor(0x90, 0xC8, 0xE8)  # soft blue
COLOR_CARD_BG   = RGBColor(0x1A, 0x2D, 0x44)  # card background

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# =============================================================================
# Environment
# =============================================================================

def find_and_load_env():
    current_dir = Path(__file__).parent
    for path in [
        current_dir / ".env",
        Path.home() / ".claude" / "skills" / "ppt-generator" / ".env",
    ]:
        if path.exists():
            load_dotenv(path, override=True)
            print(f"Loaded env: {path}")
            return
    load_dotenv(override=True)
    print("Warning: No .env file found")


# =============================================================================
# Qwen Image Generation (illustrations only)
# =============================================================================

def generate_illustration(prompt: str, out_path: str) -> Optional[str]:
    """Generate a square illustration with Qwen wanx."""
    try:
        import dashscope
        import requests
        from dashscope import ImageSynthesis

        dashscope.api_key = os.environ.get("DASHSCOPE_API_KEY", "")
        if not dashscope.api_key:
            return None

        rsp = ImageSynthesis.call(
            model="wanx2.1-t2i-turbo",
            prompt=prompt,
            n=1,
            size="1024*1024",
        )
        if rsp.status_code != 200 or not rsp.output.results:
            msg = getattr(rsp.output, "message", "unknown")
            print(f"    Illustration API error: {msg}")
            return None

        url = rsp.output.results[0].url
        data = requests.get(url, timeout=60).content
        with open(out_path, "wb") as f:
            f.write(data)
        return out_path
    except Exception as e:
        print(f"    Illustration skipped: {e}")
        return None


def make_illust_prompt(slide_info: dict) -> str:
    content = slide_info["content"][:120]
    return (
        f"Minimalist flat illustration, dark navy background, "
        f"cyan accent color, depicting: {content}. "
        f"Clean vector style, no text, professional infographic aesthetic."
    )


# =============================================================================
# Slide helpers
# =============================================================================

def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.08), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()


def add_divider(slide, x, y, w):
    line = slide.shapes.add_shape(1, x, y, w, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()


def textbox(slide, x, y, w, h, text, size, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLOR_HEADING
    return tf


def add_bullets(slide, x, y, w, h, bullets, size=20):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(size)
        run.font.color.rgb = COLOR_BODY
    return tf


# =============================================================================
# Slide builders
# =============================================================================

def build_cover(slide, slide_info):
    add_bg(slide)
    add_accent_bar(slide)

    lines = [l.strip() for l in slide_info["content"].strip().splitlines() if l.strip()]
    title    = lines[0].replace("标题：", "") if lines else "Title"
    subtitle = lines[1].replace("副标题：", "") if len(lines) > 1 else ""
    date     = lines[2].replace("时间：", "") if len(lines) > 2 else ""

    textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(2.0),
            title, 48, bold=True, align=PP_ALIGN.LEFT)
    if subtitle:
        textbox(slide, Inches(1), Inches(4.2), Inches(10), Inches(1.0),
                subtitle, 22, color=COLOR_COVER_SUB)
    if date:
        textbox(slide, Inches(1), Inches(5.4), Inches(6), Inches(0.5),
                date, 14, color=COLOR_BODY)


def build_content(slide, slide_info, illust_path=None):
    add_bg(slide)
    add_accent_bar(slide)

    lines = [l.strip() for l in slide_info["content"].strip().splitlines() if l.strip()]
    title   = lines[0]
    bullets = lines[1:]

    has_img = illust_path and os.path.exists(illust_path)
    text_w  = Inches(7.8) if has_img else Inches(12.3)

    textbox(slide, Inches(0.5), Inches(0.35), text_w, Inches(1.0),
            title, 32, bold=True)
    add_divider(slide, Inches(0.5), Inches(1.45), text_w)
    add_bullets(slide, Inches(0.5), Inches(1.65), text_w, Inches(5.6), bullets, size=21)

    if has_img:
        slide.shapes.add_picture(illust_path, Inches(8.5), Inches(1.1), Inches(4.5), Inches(4.5))


def build_data(slide, slide_info, illust_path=None):
    add_bg(slide)
    add_accent_bar(slide)

    lines = [l.strip() for l in slide_info["content"].strip().splitlines() if l.strip()]
    title = lines[0]
    items = lines[1:]

    textbox(slide, Inches(0.5), Inches(0.35), Inches(12.5), Inches(1.0),
            title, 32, bold=True)
    add_divider(slide, Inches(0.5), Inches(1.45), Inches(12.5))

    cols = 3
    card_w = Inches(3.9)
    card_h = Inches(1.6)

    for i, item in enumerate(items[:6]):
        col = i % cols
        row = i // cols
        x = Inches(0.5) + col * Inches(4.3)
        y = Inches(1.7) + row * Inches(2.0)

        card = slide.shapes.add_shape(1, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_ACCENT
        card.line.width = Pt(1)

        parts = item.split("：", 1) if "：" in item else item.split(":", 1)
        tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15),
                                       card_w - Inches(0.4), card_h - Inches(0.25))
        tf = tb.text_frame
        tf.word_wrap = True

        if len(parts) == 2:
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = parts[0].strip()
            r.font.size = Pt(13)
            r.font.color.rgb = COLOR_ACCENT
            r.font.bold = True

            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = parts[1].strip()
            r2.font.size = Pt(19)
            r2.font.color.rgb = COLOR_HEADING
            r2.font.bold = True
        else:
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = item
            r.font.size = Pt(17)
            r.font.color.rgb = COLOR_BODY


# =============================================================================
# Main
# =============================================================================

def main():
    find_and_load_env()

    parser = argparse.ArgumentParser(description="PPT Generator - text slides + AI illustrations")
    parser.add_argument("--plan", required=True, help="slides_plan.json path")
    parser.add_argument("--output", default=None, help="output directory")
    parser.add_argument("--no-illustrations", action="store_true", help="skip AI illustration generation")
    args = parser.parse_args()

    with open(args.plan, encoding="utf-8") as f:
        plan = json.load(f)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_dir = args.output or f"{OUTPUT_BASE_DIR}/{timestamp}"
    img_dir = os.path.join(out_dir, "illustrations")
    os.makedirs(img_dir, exist_ok=True)

    print(f"{'='*60}")
    print(f"PPT: {plan.get('title', 'Untitled')}")
    print(f"Slides: {len(plan['slides'])}  |  Output: {out_dir}")
    print(f"{'='*60}\n")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for slide_info in plan["slides"]:
        num   = slide_info["slide_number"]
        ptype = slide_info.get("page_type", "content")
        print(f"Slide {num} ({ptype})...")

        illust_path = None
        if not args.no_illustrations and ptype != "cover":
            illust_file = os.path.join(img_dir, f"illust-{num:02d}.png")
            print("  Generating illustration...")
            illust_path = generate_illustration(make_illust_prompt(slide_info), illust_file)
            if illust_path:
                print(f"  Illustration saved: {illust_file}")

        slide = prs.slides.add_slide(blank)

        if ptype == "cover":
            build_cover(slide, slide_info)
        elif ptype == "data":
            build_data(slide, slide_info, illust_path)
        else:
            build_content(slide, slide_info, illust_path)

        print(f"  Slide {num} done\n")

    title = plan.get("title", "presentation")
    pptx_path = os.path.join(out_dir, f"{title}.pptx")
    prs.save(pptx_path)

    print(f"{'='*60}")
    print(f"Saved: {pptx_path}")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
